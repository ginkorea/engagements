from engagement_calendar import CalendarPlotter
from map import MapPlotter
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from datetime import datetime


class EngagementsPlotter:
    def __init__(self, engagements_file='data/engagements.csv', bases_file='data/bases.csv',
                 calendar_output='calendar.png', map_output='map.png', ppt_output='engagements.pptx'):
        self.calendar_plotter = CalendarPlotter(engagements_file, calendar_output)
        self.map_plotter = MapPlotter(engagements_file, bases_file, map_output)
        self.ppt_output = ppt_output
        self.slide = None

    @staticmethod
    def get_fiscal_week():
        today = datetime.today()
        fiscal_start = datetime(today.year if today.month >= 10 else today.year - 1, 10, 1)
        fiscal_week = (today - fiscal_start).days // 7 + 1
        return fiscal_week

    def add_title(self):
        fiscal_week = self.get_fiscal_week()
        title = f"2ID/RUCD Current Engagements (Week {fiscal_week} to {fiscal_week + 4})"
        title_box = self.slide.shapes.add_textbox(0, 0, width=Inches(10), height=Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(24)
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Set no padding or margin for the title
        for paragraph in title_frame.paragraphs:
            paragraph.margin_top = 0
            paragraph.margin_bottom = 0

    def create_ppt(self):
        # Create a PowerPoint presentation
        prs = Presentation()

        # Add a slide with title and content layout
        slide_layout = prs.slide_layouts[5]  # Use a blank slide layout
        self.slide = prs.slides.add_slide(slide_layout)

        # Add a title to the slide
        self.add_title()

        # Add the calendar image at the bottom left corner
        self.slide.shapes.add_picture(self.calendar_plotter.output_image, Inches(0),
                                      prs.slide_height - Inches(5), width=Inches(5), height=Inches(5))

        # Add the map image on the right half
        self.slide.shapes.add_picture(self.map_plotter.output_image, Inches(5), prs.slide_height - Inches(6),
                                      width=Inches(5), height=Inches(6))

        # Save the PowerPoint presentation
        prs.save(self.ppt_output)

    def plot_and_save_all(self):
        self.calendar_plotter.draw_frame()
        category_counters = self.calendar_plotter.plot_engagements()
        self.calendar_plotter.save_calendar()
        self.map_plotter.plot_map()
        self.create_ppt()
        return category_counters


# Example usage
plotter = EngagementsPlotter()
category_counts = plotter.plot_and_save_all()
print(category_counts)
