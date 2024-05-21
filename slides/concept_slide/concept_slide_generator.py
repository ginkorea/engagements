from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import pandas as pd
from data.server import data_server


class ConceptSlideGenerator:
    def __init__(self, engagements_file='data/sample_engagement.csv', output_ppt='concept_slides.pptx'):
        self.engagements_file = engagements_file
        self.output_ppt = output_ppt
        self.engagements = None
        self.load_engagements()

    def load_engagements(self):
        # Load engagement data using data_server
        self.engagements = data_server.sample_engagements

    @staticmethod
    def create_shadow(element):
        shadow = element.shadow
        shadow.inherit = False
        shadow.blur_radius = Pt(4)
        shadow.distance = Pt(4)
        shadow.direction = 270  # Bottom right

    def create_slide(self, slide, engagement):
        slide_title = slide.shapes.title
        slide_title.text = engagement['Event']
        slide_title.text_frame.paragraphs[0].font.size = Pt(24)
        slide_title.top = Inches(0)
        slide_title.height = Inches(0.75)
        slide_title.left = Inches(0)
        slide_title.width = Inches(10)
        slide_title.alignment = PP_ALIGN.CENTER

        def add_box(x, y, width, height, title_text):
            box = slide.shapes.add_textbox(x, y, width, height)
            box.line.color.rgb = RGBColor(0, 0, 0)
            box.line.width = Pt(1)
            self.create_shadow(box)
            frame = box.text_frame
            frame.word_wrap = True

            title_box = slide.shapes.add_textbox(x, y, width, Inches(0.5))
            title_frame = title_box.text_frame
            title_frame.word_wrap = True

            title_p = title_frame.paragraphs[0]
            title_p.text = title_text
            title_p.font.size = Pt(12)
            title_p.font.bold = True
            title_p.font.color.rgb = RGBColor(255, 255, 255)
            title_p.font.small_caps = True
            title_p.alignment = PP_ALIGN.CENTER

            fill = title_box.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0, 0, 0)

            return frame

        left_frame = add_box(Inches(0.25), Inches(1), Inches(4.75), Inches(6), "Event Details")

        def add_paragraph(frame, label, text):
            p = frame.add_paragraph()
            p.text = f"{label}: {text}"
            p.font.size = Pt(11)
            p.alignment = PP_ALIGN.LEFT
            if label:
                p.font.bold = True

        add_paragraph(left_frame, "Date/Time", engagement['Date/Time'])
        add_paragraph(left_frame, "Location", engagement['Location'])
        add_paragraph(left_frame, "Line of Effort", engagement['Line of Effort'])
        add_paragraph(left_frame, "Subordinate Objective", engagement['Subordinate Objective'])
        add_paragraph(left_frame, "Purpose", engagement['Purpose'])
        add_paragraph(left_frame, "2ID/RUCD Attendees", engagement['2ID/RUCD Attendees'])
        add_paragraph(left_frame, "Relevant Attendees", engagement['Relevant Attendees'])
        add_paragraph(left_frame, "Uniform", engagement['Uniform'])

        # Add picture box
        img = slide.shapes.add_picture('images/img.png', Inches(5.25), Inches(1), Inches(4.5), Inches(3))
        self.create_shadow(img)

        # Add Top Line Messages / Desired Effects box
        messages_frame = add_box(Inches(5.25), Inches(4.25), Inches(4.5), Inches(4),
                                 "Top Line Messages / Desired Effects")

        for message in engagement['Messages / Effects'].split(';'):
            p = messages_frame.add_paragraph()
            p.text = message.strip()
            p.font.size = Pt(11)

        if len(messages_frame.paragraphs) < 4:
            i = len(messages_frame.paragraphs)
            while i < 4:
                p = messages_frame.add_paragraph()
                p.text = "\n"
                i += 1

        # Add Requirements box
        requirements_frame = add_box(Inches(5.25), Inches(6), Inches(4.5), Inches(1.5), "Requirements")

        def add_requirement(label, value):
            p = requirements_frame.add_paragraph()
            p.text = f"{label}: {value}"
            p.font.size = Pt(11)
            p.font.bold = True

        add_requirement("Legal Review", engagement['Legal Review'])
        add_requirement("Gift", engagement['Gift'])
        add_requirement("PAO Support", engagement['PAO Support'])
        add_requirement("Inclement Weather", engagement['Inclement Weather'])

    def generate(self):
        prs = Presentation()
        layout = prs.slide_layouts[5]  # Blank layout

        for idx, engagement in self.engagements.iterrows():
            slide = prs.slides.add_slide(layout)
            self.create_slide(slide, engagement)

        prs.save(self.output_ppt)


# Example usage
generator = ConceptSlideGenerator()
generator.generate()
